from pptx import Presentation
import tkinter as tk
from tkinter import filedialog
import io

def merge_ppts(ppt_files, output_file):
    merged_ppt = Presentation(ppt_files[0])  # Start with the first presentation
    
    for ppt_file in ppt_files[1:]:  # Skip the first file as it's already loaded
        ppt = Presentation(ppt_file)
        for slide in ppt.slides:
            new_slide = merged_ppt.slides.add_slide(merged_ppt.slide_layouts[6])  # Use a fully blank slide
            
            # Remove any existing placeholders in the new slide
            for shape in new_slide.shapes:
                if shape.is_placeholder:
                    sp = shape
                    new_slide.shapes._spTree.remove(sp._element)
            
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():  # Ensure it's a valid text box
                    text_box = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = shape.text_frame.word_wrap  # Retain text wrapping

                    for paragraph in shape.text_frame.paragraphs:
                        new_paragraph = text_frame.add_paragraph()
                        new_paragraph.alignment = paragraph.alignment  # Retain text alignment

                        for run in paragraph.runs:
                            new_run = new_paragraph.add_run()
                            new_run.text = run.text

                            # Preserve font styles
                            if run.font:
                                new_run.font.bold = True if run.font.bold else None
                                new_run.font.italic = True if run.font.italic else None
                                new_run.font.size = run.font.size if run.font.size else None
                                new_run.font.name = run.font.name if run.font.name else "Baskerville Old Face"  # Fallback

                    #'''
                elif shape.shape_type == 13:  # Picture shape type
                    image_stream = io.BytesIO(shape.image.blob)  # Convert to file-like object
                    new_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
    
    merged_ppt.save(output_file)
    print(f"Merged presentation saved as {output_file}")

def select_files():
    files = filedialog.askopenfilenames(title="Select PPT files", filetypes=[("PowerPoint Files", "*.pptx")])
    if files:
        output_file = filedialog.asksaveasfilename(
            title="Save Merged PPT As", defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")]
        )
        if output_file:
            merge_ppts(files, output_file)

# UI Setup
tk_root = tk.Tk()
tk_root.withdraw()  # Hide main window
select_files()
