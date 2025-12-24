import os
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt

def create_ppt_from_txt(txt_file, ppt_file):
    try:
        with open(txt_file, "r", encoding="utf-8") as file:
            text_sections = file.read().split("###")  # Split slides by "###"

        prs = Presentation()

        for section in text_sections:
            section = section.strip()
            if not section:
                continue  # Skip empty sections

            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
            
            # Remove default placeholders if any exist
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                slide.shapes._spTree.remove(shape._element)

            # Create a text box
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(8))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True  # Enable text wrapping

            # Process the text and detect titles
            first_paragraph = True  # Track first line to detect titles properly
            for paragraph in section.split("\n"):
                paragraph = paragraph.strip()
                if not paragraph:
                    continue  # Skip empty lines

                p = text_frame.add_paragraph() if not first_paragraph else text_frame.paragraphs[0]
                first_paragraph = False  # After first use, all new text is a paragraph

                if paragraph.startswith("**") and paragraph.endswith("**"):  
                    # Title formatting
                    paragraph = paragraph[2:-2]  # Remove "**" from both ends
                    p.text = paragraph
                    p.font.name = "Times New Roman"
                    p.font.size = Pt(24)  
                    p.font.bold = True
                else:
                    # Normal text formatting
                    p.text = paragraph
                    p.font.name = "Baskerville Old Face"
                    p.font.size = Pt(18)

        prs.save(ppt_file)
        messagebox.showinfo("Success", f"PowerPoint saved as:\n{ppt_file}")

    except Exception as e:
        messagebox.showerror("Error", f"Failed to create PPT: {e}")

# Function to select a file and process it
def select_file():
    file_path = filedialog.askopenfilename(filetypes=[("Text Files", "*.txt")])
    if not file_path:
        return  # User canceled file selection
    
    output_ppt = os.path.splitext(file_path)[0] + ".pptx"  # Save with the same name but .pptx extension
    create_ppt_from_txt(file_path, output_ppt)

# GUI Setup
root = tk.Tk()
root.title("TXT to PPT Converter")
root.geometry("400x200")

label = tk.Label(root, text="Select a .txt file to convert into PowerPoint", font=("Arial", 12))
label.pack(pady=20)

btn_select = tk.Button(root, text="Select File", command=select_file, font=("Arial", 12), bg="lightblue")
btn_select.pack(pady=10)

root.mainloop()
